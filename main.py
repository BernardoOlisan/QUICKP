from kivy.uix.floatlayout import FloatLayout
from kivy.properties import ObjectProperty
from kivymd.uix.button import MDFlatButton
from kivymd.uix.dialog import MDDialog
from email.message import EmailMessage
from kivy.core.window import Window
from pptx import Presentation
from bs4 import BeautifulSoup
from kivymd.app import MDApp
import smtplib
import wikipedia
import requests
import random
import pyttsx3
import os

class Home(FloatLayout):
    # internet function
    def see_internet(self):
        try:
            _ = requests.get('http://www.google.com/', timeout=5)
            print('Internet Connected')
        except requests.ConnectionError:
            print("No internet connection available.")

# Main code
class Manage(MDApp):
    dialog = None
    dialog2 = None
    dialog3 = None
    Window.size = (310, 550)
    title = 'QUICKP'

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.prs = Presentation(random.choice(['themes/theme1.pptx', 'themes/theme2.pptx',
                                               'themes/theme3.pptx', 'themes/theme4.pptx',
                                               'themes/theme5.pptx', 'themes/theme6.pptx',
                                               'themes/theme7.pptx', 'themes/theme8.pptx',
                                               'themes/theme9.pptx', 'themes/theme10.pptx',
                                               'themes/theme11.pptx', 'themes/theme12.pptx',
                                               'themes/theme13.pptx', 'themes/theme14.pptx',
                                               'themes/theme15.pptx'
                                               ]))

        self.engine = pyttsx3.init('sapi5')
        self.voices = self.engine.getProperty('voices')
        self.engine.setProperty('rate', 170)
        self.engine.setProperty('voice', self.voices[2].id)

        self.words_range = 380

        self.bugs = ['[1]', '[2]', '[3]', '[4]', '[5]',
                     '[6]', '[7]', '[8]', '[9]', '[10]',
                     '[11]', '[12]', '[13]', '[14]', '[15]',
                     '[16]', '[17]', '[18]', '[19]', '[20]',
                     '[21]', '[22]', '[23]', '[24]', '[25]',
                     '[26]', '[27]', '[28]', '[29]', '[30]',
                     '[31]', '[32]', '[33]', '[34]', '[35]',
                     '[36]', '[37]', '[38]', '[39]', '[40]',
                     '[41]', '[42]', '[43]', '[44]', '[45]',
                     '[46]', '[47]', '[48]', '[49]', '[50]',
                     ]

    def build(self):
        return Home()

    topic = ObjectProperty(None)
    name_file = ObjectProperty(None)

    def translate_es(self):
        wikipedia.set_lang('es')
        self.root.ids.topic.hint_text = 'Ingresa un tema'
        self.root.ids.name_file.hint_text = 'Nombre archivo'
        self.root.ids.search_label.text = 'Buscar'
        self.root.ids.back_button.text = 'Back'
        self.root.ids.create_button.text = 'Send'
        self.root.ids.email.hint_text = 'Direccion Email'
        self.dialog = MDDialog(
            title='No se realizó',
            text='No se pudo hacer. Estos son los posibles errores que pueden ser: 1- Insertó un tema no estudiado o no lo insertó. 2- No nombraste tu archivo. 3- Tu internet no funciona bien. 4- Insertaste un Email inválido',
            size_hint=(.8, .1),
            radius=[20, 20, 20, 20],
            buttons=[
                MDFlatButton(
                    text="Otra vez", text_color=self.theme_cls.primary_color,
                ),
            ],
        )
        self.dialog2 = MDDialog(
            title='AYUDA',
            text='Este software creará presentaciones de powerpoint de todos los temas, solo tiene que poner un tema y el nombre de cómo desea llamar el archivo. ya creado, se enviará al mail que introdujo.',
            size_hint=(.8, .1),
            radius=[20, 20, 20, 20],
            buttons=[
                MDFlatButton(
                    text="OK", text_color=self.theme_cls.primary_color,
                ),
            ],
        )
    def translate_en(self):
        wikipedia.set_lang('en')
        self.root.ids.topic.hint_text = 'Enter a topic'
        self.root.ids.name_file.hint_text = 'Project Name'
        self.root.ids.search_label.text = 'Search'
        self.root.ids.back_button.text = 'Back'
        self.root.ids.create_button.text = 'Send'
        self.root.ids.email.hint_text = 'Email address'
        self.dialog = MDDialog(
            title='Couldn´t make it...',
            text='Couldn´t be performed. These are the possible errors that can be: 1- You inserted a subject not studied or you didn´t insert it. 2- You didn´t name your file. 3- Your internet is not working well. 4-You insert an invalid Mail',
            size_hint=(.8, .1),
            radius=[20, 20, 20, 20],
            buttons=[
                MDFlatButton(
                    text="Try Again", text_color=self.theme_cls.primary_color,
                ),
            ],
        )
        self.dialog2 = MDDialog(
            title='HELP',
            text='This software will create powerpoint presentations of all the themes, you just have to put a theme and the name of how you want to name the file. already created the app is going to send you an email with the file.',
            size_hint=(.8, .1),
            radius=[20, 20, 20, 20],
            buttons=[
                MDFlatButton(
                    text="OK", text_color=self.theme_cls.primary_color,
                ),
            ],
        )


    Home().see_internet()


    def help_dialog(self):
        if not self.dialog2:
            self.dialog2 = MDDialog(
                title='HELP',
                text='This software will create powerpoint presentations of all the themes, you just have to put a theme and the name of how you want to name the file. already created the app is going to send you an email with the file.',
                size_hint=(.8, .1),
                radius=[20, 20, 20, 20],
                buttons=[
                    MDFlatButton(
                        text="OK", text_color=self.theme_cls.primary_color,
                    ),
                ],
            )
        self.dialog2.open()


    def Back(self):
        self.root.ids.search_button.opacity = 1
        self.root.ids.search_button.disabled = False
        self.root.ids.search_label.opacity = 1
        self.root.ids.topic.opacity = 1
        self.root.ids.topic.disabled = False
        self.root.ids.name_file.opacity = 1
        self.root.ids.name_file.disabled = False
        self.root.ids.create_button.opacity = 0
        self.root.ids.create_button.disabled = True
        self.root.ids.back_button.opacity = 0
        self.root.ids.back_button.disabled = True
        self.root.ids.send_label.opacity = 0
        self.root.ids.email.opacity = 1
        self.root.ids.email.disabled = False

    def Getting_data(self):
        def connected_to_internet(url='http://www.google.com/', timeout=5):
            try:
                _ = requests.get(url, timeout=timeout)
                print('Connected')

                self.engine.say('Searching ' + self.root.ids.topic.text + ', please wait')
                self.engine.runAndWait()

                self.root.ids.connect.opacity = 1
                self.root.ids.disconnect.opacity = 0
                self.root.ids.internet_label.opacity = 0
                try:
                    self.root.ids.search_button.opacity = 0
                    self.root.ids.search_button.disabled = True
                    self.root.ids.search_label.opacity = 0

                    # Wikipedia API System
                    self.url = wikipedia.page(self.root.ids.topic.text).url

                    # Making the web scrapper with bs4
                    self.source = requests.get(self.url)
                    self.soup = BeautifulSoup(self.source.content, features="lxml")

                    # First data
                    self.intro = wikipedia.summary(self.root.ids.topic.text, sentences=3)

                    # Second data
                    self.data_2 = self.soup.find("span", {'class': 'mw-headline'})

                    # Third data
                    self.data_3 = self.data_2.find_next("span", {'class': 'mw-headline'})

                    # Four data
                    self.data_4 = self.data_3.find_next("span", {'class': 'mw-headline'})

                    # Five data
                    self.data_5 = self.data_4.find_next("span", {'class': 'mw-headline'})

                    # Six data
                    self.data_6 = self.data_5.find_next("span", {'class': 'mw-headline'})

                    # Defining the Text of the page
                    self.title_2 = self.data_2.text
                    self.title_3 = self.data_3.text
                    self.title_4 = self.data_4.text
                    self.title_5 = self.data_5.text
                    self.title_6 = self.data_6.text


                    # Wikipedia API subtitles and split
                    for word in self.bugs:
                        if word in self.intro:
                            self.intro = self.intro.replace(word, "")

                    self.s = wikipedia.page(self.root.ids.topic.text).section(self.title_2)
                    self.second_data = self.s.split('.')[0]
                    if len(self.second_data) < self.words_range:
                        self.second_data = self.s.split('.')[0:2]
                        self.second_data1 = (".".join(self.second_data))
                    elif len(self.second_data) > self.words_range:
                        self.second_data = self.s.split('.')[0]
                        self.second_data1 = self.second_data

                    for word in self.bugs:
                        if word in self.second_data1:
                            self.second_data1 = self.second_data1.replace(word, "")

                    self.t = wikipedia.page(self.root.ids.topic.text).section(self.title_3)
                    self.third_data = self.t.split('.')[0]
                    if len(self.third_data) < self.words_range:
                        self.third_data = self.t.split('.')[0:2]
                        self.third_data1 = (".".join(self.third_data))
                    elif len(self.third_data) > self.words_range:
                        self.third_data = self.t.split('.')[0]
                        self.third_data1 = self.third_data

                    for word in self.bugs:
                        if word in self.third_data1:
                            self.third_data1 = self.third_data1.replace(word, "")

                    self.fo = wikipedia.page(self.root.ids.topic.text).section(self.title_4)
                    self.four_data = self.fo.split('.')[0]
                    if len(self.four_data) < self.words_range:
                        self.four_data = self.fo.split('.')[0:2]
                        self.four_data1 = (".".join(self.four_data))
                    elif len(self.four_data) > self.words_range:
                        self.four_data = self.fo.split('.')[0]
                        self.four_data1 = self.four_data

                    for word in self.bugs:
                        if word in self.four_data1:
                            self.four_data1 = self.four_data1.replace(word, "")

                    self.f = wikipedia.page(self.root.ids.topic.text).section(self.title_5)
                    self.five_data = self.f.split('.')[0]
                    if len(self.five_data) < self.words_range:
                        self.five_data = self.f.split('.')[0:2]
                        self.five_data1 = (".".join(self.five_data))
                    elif len(self.five_data) > self.words_range:
                        self.five_data = self.f.split('.')[0]
                        self.five_data1 = self.five_data

                    for word in self.bugs:
                        if word in self.five_data1:
                            self.five_data1 = self.five_data1.replace(word, "")

                    self.s = wikipedia.page(self.root.ids.topic.text).section(self.title_6)
                    self.six_data = self.s.split('.')[0]
                    if len(self.six_data) < self.words_range:
                        self.six_data = self.s.split('.')[0:2]
                        self.six_data1 = (".".join(self.six_data))
                    elif len(self.six_data) > self.words_range:
                        self.six_data = self.s.split('.')[0]
                        self.six_data1 = self.six_data


                    for word in self.bugs:
                        if word in self.six_data1:
                            self.six_data1 = self.six_data1.replace(word, "")

                    # making the powerpoint format
                    self.slide = self.prs.slides[0]
                    self.title = self.slide.shapes.title

                    # Second Slide introduction
                    self.slide2 = self.prs.slides[1]
                    self.second_title = self.slide2.shapes.title
                    self.second = self.slide2.placeholders[1]

                    # Third Slide
                    self.slide3 = self.prs.slides[2]
                    self.third_title = self.slide3.shapes.title
                    self.third = self.slide3.placeholders[1]

                    # four Slide
                    self.slide4 = self.prs.slides[3]
                    self.four_title = self.slide4.shapes.title
                    self.four = self.slide4.placeholders[1]

                    # Five Slide
                    self.slide5 = self.prs.slides[4]
                    self.five_title = self.slide5.shapes.title
                    self.five = self.slide5.placeholders[1]

                    # Six Slide
                    self.slide6 = self.prs.slides[5]
                    self.six_title = self.slide6.shapes.title
                    self.six = self.slide6.placeholders[1]

                    # Seven slide
                    self.slide7 = self.prs.slides[6]
                    self.seven_title = self.slide7.shapes.title
                    self.seven = self.slide7.placeholders[1]


                    # apear the create button
                    self.root.ids.create_button.opacity = 1
                    self.root.ids.create_button.disabled = False
                    self.root.ids.back_button.opacity = 1
                    self.root.ids.back_button.disabled = False

                    # hiding textfileds
                    self.root.ids.topic.opacity = 0
                    self.root.ids.topic.disabled = True
                    self.root.ids.name_file.opacity = 0
                    self.root.ids.name_file.disabled = True
                    self.root.ids.email.opacity = 0
                    self.root.ids.email.disabled = True

                    # Slide text and title defining
                    self.title.text = ('"' + self.root.ids.topic.text.upper() + '"')

                    self.second_title.text = self.root.ids.topic.text.lower()
                    self.second.text = self.intro

                    self.third_title.text = self.title_2
                    self.third.text = self.second_data1

                    self.four_title.text = self.title_3
                    self.four.text = self.third_data1

                    self.five_title.text = self.title_4
                    self.five.text = self.four_data1

                    self.six_title.text = self.title_5
                    self.six.text = self.five_data1

                    self.seven_title.text = self.title_6
                    self.seven.text = self.six_data1


                except:
                    self.root.ids.search_label.opacity = 1
                    self.root.ids.search_button.opacity = 1
                    self.root.ids.search_button.disabled = False
                    if not self.dialog:
                        self.dialog = MDDialog(
                            title='Couldn´t make it...',
                            text='Couldn´t be performed. These are the possible errors that can be: 1- You inserted a subject not studied or you didn´t insert it. 2- You didn´t name your file. 3- Your internet is not working well. 4-You insert an invalid Mail',
                            size_hint=(.8, .1),
                            radius=[20, 20, 20, 20],
                            buttons=[
                                MDFlatButton(
                                    text="Try Again", text_color=self.theme_cls.primary_color,
                                ),
                            ],
                        )
                    self.dialog.open()


            except requests.ConnectionError:
                print("No internet connection available.")
                self.root.ids.search_label.opacity = 0
                self.root.ids.search_button.opacity = 0
                self.root.ids.search_button.disabled = True
                self.root.ids.disconnect.opacity = 1
                self.root.ids.internet_label.opacity = 1
                self.root.ids.connect.opacity = 0
                self.root.ids.internet_button.opacity = 1
                self.root.ids.internet_button.disabled = False

        connected_to_internet()

    def try_again(self):
        self.root.ids.search_button.opacity = 1
        self.root.ids.search_button.disabled = False
        self.root.ids.search_label.opacity = 1
        self.root.ids.internet_button.opacity = 0
        self.root.ids.internet_button.disabled = True
        self.root.ids.internet_label.opacity = 0


    def send(self):
        try:
            self.engine.say('We are sending it, hold on')
            self.engine.runAndWait()

            self.EMAIL_ADDRESS = 'bernardoolisan@gmail.com'                           # Here you enter your email dir
            self.EMAIL_PASSWORD = 'gabriela26'                                    # here your email password

            self.msg = EmailMessage()
            self.msg['Subject'] = 'Here is your Presentation of ' + self.root.ids.topic.text
            self.msg['From'] = self.EMAIL_ADDRESS
            self.msg['To'] = self.root.ids.email.text
            self.msg.set_content('Thanks for using QUICKP\nYour Presentation... ')

            def save(name):
                self.prs.save(name + '.pptx')

            save(self.root.ids.name_file.text)

            with open(self.root.ids.name_file.text + '.pptx', 'rb') as f:
                self.file_data = f.read()
                self.file_name = f.name

            self.msg.add_attachment(self.file_data, maintype='application', subtype='octed-stream', filename=self.file_name)

            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
                smtp.login(self.EMAIL_ADDRESS, self.EMAIL_PASSWORD)
                smtp.send_message(self.msg)

                print('Successfully sended!')
                self.root.ids.send_label.opacity = 1

                self.engine.say('Successfully sended!')
                self.engine.runAndWait()

            os.remove(self.file_name)

        except:
            if not self.dialog:
                self.dialog = MDDialog(
                    title='Couldn´t make it...',
                    text='Couldn´t be performed. These are the possible errors that can be: 1- You inserted a subject not studied or you didn´t insert it. 2- You didn´t name your file. 3- Your internet is not working well. 4- You insert an invalid Mail',
                    size_hint=(.8, .1),
                    radius=[20, 20, 20, 20],
                    buttons=[
                        MDFlatButton(
                            text="Try Again", text_color=self.theme_cls.primary_color,
                        ),
                    ],
                )
            self.dialog.open()
            os.remove(self.file_name)


if __name__ == '__main__':
    Manage().run()
